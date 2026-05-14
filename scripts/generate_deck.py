"""Generate the Energy Hunter pitch deck (.pptx) for the 5-minute hackathon demo.

Run:
    python scripts/generate_deck.py

Output:
    presentations/energy_hunter_pitch.pptx

The deck uses an eco-friendly palette (deep green, leaf green, sand, charcoal)
and embeds illustrative images generated with Pillow so the file is fully
self-contained and reproducible without internet access.
"""

from __future__ import annotations

import math
import random
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------- #
# Eco palette                                                                 #
# --------------------------------------------------------------------------- #

DEEP_GREEN = RGBColor(0x10, 0x4F, 0x3A)   # bosque profundo
LEAF_GREEN = RGBColor(0x2E, 0x8B, 0x57)   # hoja
LIME = RGBColor(0x9A, 0xCD, 0x32)         # acento brillante
SAND = RGBColor(0xF5, 0xEF, 0xE0)         # fondo cálido neutro
CHARCOAL = RGBColor(0x1F, 0x2A, 0x24)     # texto principal
SLATE = RGBColor(0x4B, 0x5A, 0x52)        # texto secundario
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ALERT = RGBColor(0xC2, 0x4D, 0x2C)        # ámbar tierra para alertas

# Pillow-friendly tuples
P_DEEP = (0x10, 0x4F, 0x3A)
P_LEAF = (0x2E, 0x8B, 0x57)
P_LIME = (0x9A, 0xCD, 0x32)
P_SAND = (0xF5, 0xEF, 0xE0)
P_CHARCOAL = (0x1F, 0x2A, 0x24)
P_WHITE = (0xFF, 0xFF, 0xFF)
P_ALERT = (0xC2, 0x4D, 0x2C)
P_SKY = (0xCDE7DA)


# --------------------------------------------------------------------------- #
# Paths                                                                       #
# --------------------------------------------------------------------------- #

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "presentations" / "assets"
OUTPUT = ROOT / "presentations" / "energy_hunter_pitch.pptx"
ASSETS.mkdir(parents=True, exist_ok=True)


# --------------------------------------------------------------------------- #
# Image generation helpers (Pillow)                                           #
# --------------------------------------------------------------------------- #

def _font(size: int, bold: bool = False) -> ImageFont.ImageFont:
    """Best-effort font loader; falls back to default bitmap font on Windows."""
    candidates = (
        ["arialbd.ttf", "calibrib.ttf", "segoeuib.ttf"]
        if bold
        else ["arial.ttf", "calibri.ttf", "segoeui.ttf"]
    )
    for name in candidates:
        try:
            return ImageFont.truetype(name, size)
        except OSError:
            continue
    return ImageFont.load_default()


def _gradient(size: tuple[int, int], top: tuple[int, int, int], bottom: tuple[int, int, int]) -> Image.Image:
    w, h = size
    base = Image.new("RGB", size, top)
    for y in range(h):
        t = y / max(h - 1, 1)
        r = int(top[0] * (1 - t) + bottom[0] * t)
        g = int(top[1] * (1 - t) + bottom[1] * t)
        b = int(top[2] * (1 - t) + bottom[2] * t)
        ImageDraw.Draw(base).line([(0, y), (w, y)], fill=(r, g, b))
    return base


def make_hero_image(path: Path) -> None:
    """Cover hero: stylised city skyline with a sun/leaf and energy lines."""
    W, H = 1600, 900
    img = _gradient((W, H), (0xE8, 0xF1, 0xE3), (0xBF, 0xDC, 0xC4))
    draw = ImageDraw.Draw(img, "RGBA")

    # Sun / leaf circle
    cx, cy, r = 1280, 240, 150
    draw.ellipse((cx - r, cy - r, cx + r, cy + r), fill=(0xCB, 0xE6, 0x8E, 255))
    # Leaf vein
    draw.line((cx - r + 30, cy + r - 30, cx + r - 30, cy - r + 30), fill=P_DEEP, width=8)

    # Skyline buildings
    random.seed(7)
    base_y = 720
    x = 80
    while x < W - 80:
        bw = random.randint(70, 160)
        bh = random.randint(180, 420)
        color = random.choice([P_DEEP, P_LEAF, (0x1A, 0x6E, 0x52), (0x14, 0x5A, 0x42)])
        draw.rectangle((x, base_y - bh, x + bw, base_y), fill=color)
        # Windows grid
        for wy in range(base_y - bh + 20, base_y - 20, 28):
            for wx in range(x + 12, x + bw - 12, 22):
                if random.random() > 0.35:
                    draw.rectangle((wx, wy, wx + 10, wy + 14), fill=(0xF5, 0xEF, 0xE0, 220))
        x += bw + 12

    # Ground line
    draw.rectangle((0, base_y, W, H), fill=(0x10, 0x3A, 0x2C))

    # Energy waves
    for i, amp in enumerate((40, 28, 18)):
        pts = []
        for px in range(0, W, 8):
            py = 820 + int(math.sin((px / 80) + i) * amp * 0.4)
            pts.append((px, py))
        draw.line(pts, fill=(0x9A, 0xCD, 0x32, 220 - i * 40), width=4)

    img.save(path, "PNG", optimize=True)


def make_problem_image(path: Path) -> None:
    """Stylised meter dial with overload zone."""
    W, H = 1200, 800
    img = Image.new("RGB", (W, H), P_SAND)
    draw = ImageDraw.Draw(img, "RGBA")

    cx, cy, r = W // 2, H // 2 + 40, 280
    # Dial background
    draw.ellipse((cx - r, cy - r, cx + r, cy + r), fill=P_WHITE, outline=P_CHARCOAL, width=6)

    # Coloured arcs: green / amber / red
    draw.arc((cx - r + 20, cy - r + 20, cx + r - 20, cy + r - 20),
             start=180, end=240, fill=P_LEAF, width=40)
    draw.arc((cx - r + 20, cy - r + 20, cx + r - 20, cy + r - 20),
             start=240, end=300, fill=(0xD4, 0xA0, 0x2A), width=40)
    draw.arc((cx - r + 20, cy - r + 20, cx + r - 20, cy + r - 20),
             start=300, end=360, fill=P_ALERT, width=40)

    # Needle pointing into the red zone
    angle = math.radians(335)
    nx = cx + int(math.cos(angle) * (r - 60))
    ny = cy + int(math.sin(angle) * (r - 60))
    draw.line((cx, cy, nx, ny), fill=P_CHARCOAL, width=10)
    draw.ellipse((cx - 14, cy - 14, cx + 14, cy + 14), fill=P_CHARCOAL)

    # Label
    f = _font(48, bold=True)
    draw.text((cx - 220, cy + r + 20), "Consumo fuera de control", fill=P_CHARCOAL, font=f)

    img.save(path, "PNG", optimize=True)


def make_dashboard_mock(path: Path) -> None:
    """Stylised dashboard mock: map + KPIs + chart."""
    W, H = 1600, 900
    img = Image.new("RGB", (W, H), P_SAND)
    draw = ImageDraw.Draw(img, "RGBA")

    # Top bar
    draw.rectangle((0, 0, W, 80), fill=P_DEEP)
    draw.text((40, 22), "Energy Hunter — Panel B2B", fill=P_WHITE, font=_font(32, bold=True))
    draw.ellipse((W - 60, 28, W - 36, 52), fill=P_LIME)

    # Sidebar (building list)
    draw.rectangle((0, 80, 320, H), fill=(0xEC, 0xE3, 0xCE))
    draw.text((24, 100), "Edificios", fill=P_CHARCOAL, font=_font(24, bold=True))
    for i, name in enumerate(["Sede Madrid", "Planta Sevilla", "Centro Bilbao", "Oficina Valencia"]):
        y = 150 + i * 60
        draw.rectangle((16, y, 304, y + 48), fill=P_WHITE, outline=P_LEAF, width=2)
        draw.ellipse((30, y + 16, 46, y + 32), fill=P_LEAF if i == 0 else (0xCC, 0xCC, 0xCC))
        draw.text((58, y + 14), name, fill=P_CHARCOAL, font=_font(20))

    # Map area
    map_x0, map_y0, map_x1, map_y1 = 340, 100, 1000, 520
    draw.rectangle((map_x0, map_y0, map_x1, map_y1), fill=(0xDC, 0xEA, 0xD8), outline=P_LEAF, width=3)
    # Fake map grid
    for gx in range(map_x0, map_x1, 40):
        draw.line((gx, map_y0, gx, map_y1), fill=(0xC9, 0xDC, 0xC2), width=1)
    for gy in range(map_y0, map_y1, 40):
        draw.line((map_x0, gy, map_x1, gy), fill=(0xC9, 0xDC, 0xC2), width=1)
    # Markers
    pins = [(map_x0 + 180, map_y0 + 140, P_LEAF),
            (map_x0 + 360, map_y0 + 220, P_LEAF),
            (map_x0 + 500, map_y0 + 120, P_ALERT),  # anomaly
            (map_x0 + 280, map_y0 + 320, P_LEAF)]
    for (px, py, color) in pins:
        draw.ellipse((px - 14, py - 14, px + 14, py + 14), fill=color, outline=P_WHITE, width=3)
        if color == P_ALERT:
            draw.polygon([(px, py - 30), (px - 12, py - 8), (px + 12, py - 8)], fill=P_ALERT)
    draw.text((map_x0 + 16, map_y0 + 12), "Mapa interactivo", fill=P_DEEP, font=_font(22, bold=True))

    # KPI cards
    kpis = [("kWh hoy", "1.842"), ("€/MWh ahora", "112"), ("Anomalías", "3"), ("Modo Eco", "2 fases")]
    for i, (label, value) in enumerate(kpis):
        x0 = 1020 + (i % 2) * 280
        y0 = 100 + (i // 2) * 200
        draw.rectangle((x0, y0, x0 + 260, y0 + 180), fill=P_WHITE, outline=P_LEAF, width=2)
        draw.text((x0 + 20, y0 + 18), label, fill=P_CHARCOAL, font=_font(20))
        draw.text((x0 + 20, y0 + 58), value, fill=P_DEEP, font=_font(60, bold=True))

    # Chart area
    cx0, cy0, cx1, cy1 = 340, 540, 1560, 860
    draw.rectangle((cx0, cy0, cx1, cy1), fill=P_WHITE, outline=P_LEAF, width=2)
    draw.text((cx0 + 16, cy0 + 12), "Consumo por fase (últimas 6 horas)", fill=P_DEEP, font=_font(22, bold=True))
    # Axes
    draw.line((cx0 + 40, cy1 - 20, cx1 - 20, cy1 - 20), fill=P_CHARCOAL, width=2)
    draw.line((cx0 + 40, cy0 + 50, cx0 + 40, cy1 - 20), fill=P_CHARCOAL, width=2)
    # Series
    random.seed(11)
    for series_color in (P_LEAF, P_DEEP, (0x4F, 0xA0, 0x70)):
        pts = []
        for k in range(40):
            x = cx0 + 50 + k * ((cx1 - cx0 - 70) // 40)
            y = cy1 - 30 - random.randint(20, 220)
            pts.append((x, y))
        draw.line(pts, fill=series_color, width=4)
    # Anomaly marker
    ax, ay = cx0 + 50 + 26 * ((cx1 - cx0 - 70) // 40), cy0 + 90
    draw.polygon([(ax, ay - 14), (ax - 12, ay + 8), (ax + 12, ay + 8)], fill=P_ALERT)
    draw.text((ax + 16, ay - 8), "Anomalía domingo", fill=P_ALERT, font=_font(18, bold=True))

    img.save(path, "PNG", optimize=True)


def make_panic_image(path: Path) -> None:
    """Big eco button + checkboxes for the panic flow."""
    W, H = 1200, 800
    img = Image.new("RGB", (W, H), P_SAND)
    draw = ImageDraw.Draw(img, "RGBA")

    # Card
    draw.rectangle((80, 80, W - 80, H - 80), fill=P_WHITE, outline=P_LEAF, width=4)
    draw.text((120, 110), "Botón de Pánico granular", fill=P_DEEP, font=_font(40, bold=True))
    draw.text((120, 170), "Selecciona qué fases pasan a Modo Eco (-40%)", fill=P_CHARCOAL, font=_font(24))

    # Checkboxes
    items = [("Fase A — Iluminación oficinas", True),
             ("Fase B — Climatización plantas 1-3", True),
             ("Fase C — Servidores críticos", False),
             ("Fase D — Cargadores VE", True)]
    for i, (text, checked) in enumerate(items):
        y = 230 + i * 70
        # Box
        draw.rectangle((140, y, 180, y + 40), fill=P_WHITE, outline=P_DEEP, width=3)
        if checked:
            draw.line((148, y + 22, 162, y + 34), fill=P_LEAF, width=6)
            draw.line((162, y + 34, 174, y + 10), fill=P_LEAF, width=6)
        draw.text((210, y + 6), text, fill=P_CHARCOAL, font=_font(26))

    # Panic button (eco green, accessible label)
    draw.rounded_rectangle((140, 580, 700, 680), radius=24, fill=P_DEEP)
    draw.text((180, 612), "Activar Modo Eco", fill=P_WHITE, font=_font(36, bold=True))

    # Result chip
    draw.rounded_rectangle((740, 580, 1080, 680), radius=24, fill=P_LIME)
    draw.text((780, 612), "−40% consumo", fill=P_CHARCOAL, font=_font(30, bold=True))

    img.save(path, "PNG", optimize=True)


def make_architecture_image(path: Path) -> None:
    """Simple architecture diagram: Streamlit ↔ Mock API ↔ Kiro stub."""
    W, H = 1400, 700
    img = Image.new("RGB", (W, H), P_SAND)
    draw = ImageDraw.Draw(img, "RGBA")

    boxes = [
        (60, 260, 320, 440, "Streamlit\nDashboard", P_LEAF),
        (440, 100, 720, 280, "Smart Breaker\nMock API", P_DEEP),
        (440, 420, 720, 600, "SIOS\nPrecios", P_DEEP),
        (840, 260, 1100, 440, "Kiro\nWebhook", P_LEAF),
        (1180, 260, 1380, 440, "Modo Eco\n−40%", P_LIME),
    ]
    for (x0, y0, x1, y1, label, color) in boxes:
        draw.rounded_rectangle((x0, y0, x1, y1), radius=18, fill=color, outline=P_CHARCOAL, width=3)
        f = _font(28, bold=True)
        text_color = P_WHITE if color != P_LIME else P_CHARCOAL
        # Center text
        lines = label.split("\n")
        total_h = len(lines) * 36
        ty = (y0 + y1) // 2 - total_h // 2
        for line in lines:
            bbox = draw.textbbox((0, 0), line, font=f)
            tw = bbox[2] - bbox[0]
            draw.text(((x0 + x1) // 2 - tw // 2, ty), line, fill=text_color, font=f)
            ty += 36

    # Arrows
    def arrow(p0, p1):
        draw.line([p0, p1], fill=P_CHARCOAL, width=4)
        ax, ay = p1
        draw.polygon([(ax, ay - 10), (ax, ay + 10), (ax + 18, ay)], fill=P_CHARCOAL)

    arrow((320, 350), (438, 190))
    arrow((320, 350), (438, 510))
    arrow((720, 190), (840, 330))
    arrow((720, 510), (840, 370))
    arrow((1100, 350), (1178, 350))

    # Title
    draw.text((60, 40), "Arquitectura local 100% reproducible", fill=P_DEEP, font=_font(34, bold=True))

    img.save(path, "PNG", optimize=True)


def make_savings_chart(path: Path) -> None:
    """Bar chart: kWh before vs after Eco Mode."""
    W, H = 1200, 720
    img = Image.new("RGB", (W, H), P_SAND)
    draw = ImageDraw.Draw(img, "RGBA")

    draw.text((60, 40), "Impacto inmediato del Modo Eco", fill=P_DEEP, font=_font(34, bold=True))

    categories = [("Iluminación", 120, 72),
                  ("Climatización", 240, 144),
                  ("Cargadores VE", 180, 108),
                  ("Total fases", 540, 324)]

    # Axes
    ox, oy = 120, 620
    draw.line((ox, oy, W - 60, oy), fill=P_CHARCOAL, width=3)
    draw.line((ox, 140, ox, oy), fill=P_CHARCOAL, width=3)

    bar_w = 90
    gap = 160
    max_val = 600
    for i, (label, before, after) in enumerate(categories):
        x = ox + 60 + i * (bar_w * 2 + gap - 90)
        # Before bar (charcoal)
        bh = int((before / max_val) * (oy - 160))
        draw.rectangle((x, oy - bh, x + bar_w, oy), fill=(0x4B, 0x5A, 0x52))
        draw.text((x, oy - bh - 30), f"{before}", fill=P_CHARCOAL, font=_font(22, bold=True))
        # After bar (lime)
        ah = int((after / max_val) * (oy - 160))
        draw.rectangle((x + bar_w + 10, oy - ah, x + bar_w * 2 + 10, oy), fill=P_LEAF)
        draw.text((x + bar_w + 10, oy - ah - 30), f"{after}", fill=P_DEEP, font=_font(22, bold=True))
        # Label
        draw.text((x - 4, oy + 12), label, fill=P_CHARCOAL, font=_font(20))

    # Legend
    draw.rectangle((W - 360, 100, W - 340, 120), fill=(0x4B, 0x5A, 0x52))
    draw.text((W - 320, 96), "Antes (kWh)", fill=P_CHARCOAL, font=_font(22))
    draw.rectangle((W - 360, 140, W - 340, 160), fill=P_LEAF)
    draw.text((W - 320, 136), "Con Modo Eco (−40%)", fill=P_CHARCOAL, font=_font(22))

    img.save(path, "PNG", optimize=True)


def make_aws_architecture_image(path: Path) -> None:
    """AWS deployment diagram: edge devices → IoT Core → processing → dashboard."""
    W, H = 1600, 800
    img = Image.new("RGB", (W, H), P_SAND)
    draw = ImageDraw.Draw(img, "RGBA")

    draw.text((50, 30), "Despliegue y escalado en AWS", fill=P_DEEP, font=_font(34, bold=True))
    draw.text((50, 78), "De edificios físicos a un panel global, sin servidores que mantener",
              fill=P_CHARCOAL, font=_font(20))

    # Three layers (columns of boxes)
    layers = [
        # (x0, label, color, items)
        (60, "Edge / Edificio", P_LEAF,
         ["Smart Breakers", "Greengrass\n(opcional)", "Sensor SIOS"]),
        (440, "AWS IoT Core", P_DEEP,
         ["MQTT Broker", "Device Shadow", "Rules Engine"]),
        (820, "Procesado", P_DEEP,
         ["Lambda\n(anomalías)", "Timestream\n(series)", "S3 + Athena\n(histórico)"]),
        (1200, "Capa cliente", P_LEAF,
         ["ECS Fargate\nStreamlit", "API Gateway\n+ Cognito", "CloudFront\n(CDN)"]),
    ]

    box_w = 320
    box_h = 110
    gap = 18
    title_y = 150
    for (x0, label, color, items) in layers:
        # Layer header pill
        draw.rounded_rectangle((x0, title_y, x0 + box_w, title_y + 50),
                               radius=12, fill=color)
        f_h = _font(22, bold=True)
        bbox = draw.textbbox((0, 0), label, font=f_h)
        tw = bbox[2] - bbox[0]
        draw.text((x0 + (box_w - tw) // 2, title_y + 12), label,
                  fill=P_WHITE, font=f_h)
        # Items
        y = title_y + 70
        for item in items:
            draw.rounded_rectangle((x0, y, x0 + box_w, y + box_h),
                                   radius=10, fill=P_WHITE,
                                   outline=color, width=3)
            f = _font(20, bold=True)
            lines = item.split("\n")
            total_h = len(lines) * 28
            ty = y + (box_h - total_h) // 2
            for line in lines:
                bbox = draw.textbbox((0, 0), line, font=f)
                tw = bbox[2] - bbox[0]
                draw.text((x0 + (box_w - tw) // 2, ty), line,
                          fill=P_CHARCOAL, font=f)
                ty += 28
            y += box_h + gap

    # Connecting arrows between layer columns (centered between boxes)
    def arrow(p0, p1):
        draw.line([p0, p1], fill=P_CHARCOAL, width=4)
        ax, ay = p1
        draw.polygon([(ax - 16, ay - 9), (ax - 16, ay + 9), (ax, ay)],
                     fill=P_CHARCOAL)

    mid_y = title_y + 70 + box_h + gap + box_h // 2  # middle row vertical center
    for i in range(3):
        x_from = 60 + (i + 1) * 380 - 60  # right edge of column i (380 = box_w + 60 gap)
        x_to = x_from + 60
        arrow((x_from, mid_y), (x_to, mid_y))

    # Footer line: scaling note
    draw.rounded_rectangle((50, 700, W - 50, 770), radius=12,
                           fill=P_LIME)
    draw.text((70, 718),
              "Escala lineal: 10 → 10.000 edificios sin tocar la arquitectura. "
              "Pago por uso, multi-región y tolerancia a fallos por defecto.",
              fill=P_CHARCOAL, font=_font(20, bold=True))

    img.save(path, "PNG", optimize=True)


def make_cost_comparison_image(path: Path) -> None:
    """Side-by-side bars: monthly AWS cost vs monthly energy savings for 3 tiers."""
    W, H = 1400, 760
    img = Image.new("RGB", (W, H), P_SAND)
    draw = ImageDraw.Draw(img, "RGBA")

    draw.text((50, 30), "Coste AWS vs ahorro energético generado",
              fill=P_DEEP, font=_font(32, bold=True))
    draw.text((50, 75),
              "Tres escenarios de cliente. Barras en €/mes. "
              "Verde = lo que pagas a AWS; Lima = lo que ahorras en factura eléctrica.",
              fill=P_CHARCOAL, font=_font(18))

    tiers = [
        ("Pyme\n3 edificios · 30 fases", 95, 720),
        ("Mediana empresa\n15 edificios · 180 fases", 380, 4200),
        ("Corporativo\n80 edificios · 1.200 fases", 1850, 26000),
    ]

    # Plot area
    ox, oy = 110, 640
    plot_h = 470
    max_val = max(saving for _, _, saving in tiers) * 1.1

    draw.line((ox, 130, ox, oy), fill=P_CHARCOAL, width=3)
    draw.line((ox, oy, W - 60, oy), fill=P_CHARCOAL, width=3)

    group_w = (W - 60 - ox) // len(tiers)
    bar_w = 110

    for i, (label, cost, saving) in enumerate(tiers):
        gx = ox + i * group_w + group_w // 2
        # AWS cost (deep green)
        ch = int((cost / max_val) * plot_h)
        draw.rectangle((gx - bar_w - 10, oy - ch, gx - 10, oy),
                       fill=P_DEEP)
        draw.text((gx - bar_w - 10, oy - ch - 32),
                  f"{cost} €", fill=P_DEEP, font=_font(22, bold=True))
        # Savings (lime/leaf)
        sh = int((saving / max_val) * plot_h)
        draw.rectangle((gx + 10, oy - sh, gx + bar_w + 10, oy),
                       fill=P_LEAF)
        draw.text((gx + 10, oy - sh - 32),
                  f"{saving} €", fill=P_DEEP, font=_font(22, bold=True))
        # Net delta chip
        net = saving - cost
        chip_w = 200
        draw.rounded_rectangle((gx - chip_w // 2, oy + 90,
                                gx + chip_w // 2, oy + 130),
                               radius=14, fill=P_LIME)
        bbox = draw.textbbox((0, 0), f"+ {net} € netos / mes",
                             font=_font(18, bold=True))
        tw = bbox[2] - bbox[0]
        draw.text((gx - tw // 2, oy + 98), f"+ {net} € netos / mes",
                  fill=P_CHARCOAL, font=_font(18, bold=True))
        # Tier label
        f = _font(18, bold=True)
        for j, line in enumerate(label.split("\n")):
            bbox = draw.textbbox((0, 0), line, font=f)
            tw = bbox[2] - bbox[0]
            draw.text((gx - tw // 2, oy + 18 + j * 26),
                      line, fill=P_CHARCOAL, font=f)

    # Legend
    draw.rectangle((W - 360, 100, W - 340, 122), fill=P_DEEP)
    draw.text((W - 320, 96), "Coste AWS / mes", fill=P_CHARCOAL,
              font=_font(20, bold=True))
    draw.rectangle((W - 360, 138, W - 340, 160), fill=P_LEAF)
    draw.text((W - 320, 134), "Ahorro factura / mes",
              fill=P_CHARCOAL, font=_font(20, bold=True))

    img.save(path, "PNG", optimize=True)


def make_aws_cost_breakdown_image(path: Path) -> None:
    """Stacked horizontal bar showing AWS cost split by service for each tier."""
    W, H = 1400, 720
    img = Image.new("RGB", (W, H), P_SAND)
    draw = ImageDraw.Draw(img, "RGBA")

    draw.text((50, 30), "Desglose del coste AWS por servicio",
              fill=P_DEEP, font=_font(32, bold=True))
    draw.text((50, 78),
              "Estimación mensual en eu-west-1 con telemetría cada 30 s.",
              fill=P_CHARCOAL, font=_font(18))

    # tier_label, total, breakdown[(service, amount, color)]
    tiers = [
        ("Pyme · 30 fases · 95 €", 95, [
            ("IoT Core", 12, P_DEEP),
            ("Timestream", 28, P_LEAF),
            ("Lambda", 6, (0x4F, 0xA0, 0x70)),
            ("Fargate", 32, (0x14, 0x5A, 0x42)),
            ("CloudFront/Cognito/S3", 17, P_LIME),
        ]),
        ("Mediana · 180 fases · 380 €", 380, [
            ("IoT Core", 60, P_DEEP),
            ("Timestream", 130, P_LEAF),
            ("Lambda", 22, (0x4F, 0xA0, 0x70)),
            ("Fargate", 110, (0x14, 0x5A, 0x42)),
            ("CloudFront/Cognito/S3", 58, P_LIME),
        ]),
        ("Corporativo · 1.200 fases · 1.850 €", 1850, [
            ("IoT Core", 320, P_DEEP),
            ("Timestream", 680, P_LEAF),
            ("Lambda", 110, (0x4F, 0xA0, 0x70)),
            ("Fargate", 480, (0x14, 0x5A, 0x42)),
            ("CloudFront/Cognito/S3", 260, P_LIME),
        ]),
    ]

    bar_x0 = 380
    bar_x1 = W - 80
    bar_full_w = bar_x1 - bar_x0
    row_h = 90
    y = 160

    f_label = _font(20, bold=True)
    f_small = _font(14, bold=True)

    for (label, total, items) in tiers:
        # Tier label
        for j, line in enumerate(label.split(" · ")):
            draw.text((50, y + j * 24), line, fill=P_CHARCOAL, font=f_label)
        # Stacked bar
        x = bar_x0
        for (service, amount, color) in items:
            seg_w = int(bar_full_w * (amount / total))
            draw.rectangle((x, y, x + seg_w, y + 60), fill=color)
            if seg_w > 80:
                draw.text((x + 8, y + 10), service,
                          fill=P_WHITE if color != P_LIME else P_CHARCOAL,
                          font=f_small)
                draw.text((x + 8, y + 32), f"{amount} €",
                          fill=P_WHITE if color != P_LIME else P_CHARCOAL,
                          font=f_small)
            x += seg_w
        # Border
        draw.rectangle((bar_x0, y, bar_x1, y + 60), outline=P_CHARCOAL, width=2)
        y += row_h + 30

    # Footer note
    draw.rounded_rectangle((50, H - 80, W - 50, H - 30),
                           radius=12, fill=P_LIME)
    draw.text((70, H - 68),
              "Pago por uso. Sin licencias por dispositivo. Sin compromiso anual.",
              fill=P_CHARCOAL, font=_font(20, bold=True))

    img.save(path, "PNG", optimize=True)


def make_roadmap_image(path: Path) -> None:
    """Horizontal roadmap with milestones."""
    W, H = 1400, 500
    img = Image.new("RGB", (W, H), P_SAND)
    draw = ImageDraw.Draw(img, "RGBA")
    draw.text((60, 40), "Roadmap", fill=P_DEEP, font=_font(36, bold=True))

    # Track
    draw.rectangle((80, 240, W - 80, 270), fill=P_LEAF)
    milestones = [
        (180, "MVP demo\nhackathon", P_LIME),
        (480, "Piloto\n3 edificios", P_LEAF),
        (780, "Integración\nbreakers reales", P_LEAF),
        (1080, "Multi-tenant\n+ alertas IA", P_DEEP),
        (1320, "Mercado\nEU", P_DEEP),
    ]
    for (x, label, color) in milestones:
        draw.ellipse((x - 22, 233, x + 22, 277), fill=color, outline=P_CHARCOAL, width=3)
        # Label box
        lines = label.split("\n")
        f = _font(22, bold=True)
        ty = 310
        for line in lines:
            bbox = draw.textbbox((0, 0), line, font=f)
            tw = bbox[2] - bbox[0]
            draw.text((x - tw // 2, ty), line, fill=P_CHARCOAL, font=f)
            ty += 28

    img.save(path, "PNG", optimize=True)


# --------------------------------------------------------------------------- #
# Slide builder helpers (python-pptx)                                         #
# --------------------------------------------------------------------------- #

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_background(slide, color: RGBColor) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.shadow.inherit = False
    # Send to back by reordering xml
    spTree = rect._element.getparent()
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)


def add_side_band(slide, color: RGBColor, width=Inches(0.35)) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, SLIDE_H)
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = CHARCOAL,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    bullets: list[str],
    *,
    size: int = 22,
    color: RGBColor = CHARCOAL,
    bullet_color: RGBColor = LEAF_GREEN,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # Bullet marker
        marker = p.add_run()
        marker.text = "▎ "
        marker.font.size = Pt(size)
        marker.font.bold = True
        marker.font.color.rgb = bullet_color
        marker.font.name = "Calibri"
        # Text
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.space_after = Pt(8)
    return tb


def add_footer(slide, page_num: int, total: int) -> None:
    add_textbox(
        slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.35),
        "Energy Hunter · Hackathon Demo", size=11, color=SLATE,
    )
    add_textbox(
        slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.35),
        f"{page_num} / {total}", size=11, color=SLATE, align=PP_ALIGN.RIGHT,
    )


def add_speaker_notes(slide, text: str) -> None:
    """Attach speaker notes to the slide (~25-30 s of speech)."""
    slide.notes_slide.notes_text_frame.text = text


# --------------------------------------------------------------------------- #
# Slide layouts                                                               #
# --------------------------------------------------------------------------- #

TOTAL_SLIDES = 13


def build_cover(prs: Presentation, hero_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_background(slide, SAND)

    # Hero image as full-bleed top band
    slide.shapes.add_picture(str(hero_path), 0, 0, width=SLIDE_W, height=Inches(4.2))

    # Green title bar overlay
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.6), SLIDE_W, Inches(0.6))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = DEEP_GREEN

    add_textbox(
        slide, Inches(0.7), Inches(3.62), Inches(12), Inches(0.6),
        "Energy Hunter", size=28, bold=True, color=WHITE,
    )

    # Subtitle block
    add_textbox(
        slide, Inches(0.7), Inches(4.5), Inches(12), Inches(0.7),
        "Gestión inteligente de energía para edificios B2B",
        size=34, bold=True, color=DEEP_GREEN,
    )
    add_textbox(
        slide, Inches(0.7), Inches(5.25), Inches(12), Inches(0.6),
        "Visibilidad en tiempo real · Detección de anomalías · Modo Eco con un clic",
        size=20, color=SLATE,
    )

    # Tagline pill
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(6.1), Inches(5.2), Inches(0.55))
    pill.line.fill.background()
    pill.fill.solid()
    pill.fill.fore_color.rgb = LIME
    add_textbox(
        slide, Inches(0.7), Inches(6.1), Inches(5.2), Inches(0.55),
        "Menos consumo · Menos emisiones · Más control",
        size=16, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    add_textbox(
        slide, Inches(0.7), Inches(6.8), Inches(12), Inches(0.4),
        "Avance MVP · Hackathon 2026", size=14, color=SLATE,
    )


def build_title_slide(prs: Presentation, page: int, eyebrow: str, title: str, subtitle: str) -> object:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, SAND)
    add_side_band(slide, DEEP_GREEN)
    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
                eyebrow.upper(), size=12, bold=True, color=LEAF_GREEN)
    add_textbox(slide, Inches(0.7), Inches(0.85), Inches(12), Inches(0.9),
                title, size=34, bold=True, color=DEEP_GREEN)
    add_textbox(slide, Inches(0.7), Inches(1.7), Inches(12), Inches(0.6),
                subtitle, size=18, color=SLATE)
    add_footer(slide, page, TOTAL_SLIDES)
    return slide


def build_problem(prs: Presentation, page: int, image: Path) -> None:
    slide = build_title_slide(
        prs, page,
        "El problema",
        "Los edificios consumen sin saber dónde ni cuándo",
        "Datos dispersos, facturas sorpresa y nula reacción ante picos de consumo",
    )
    slide.shapes.add_picture(str(image), Inches(7.6), Inches(2.4), width=Inches(5.2))
    add_bullets(
        slide, Inches(0.7), Inches(2.6), Inches(6.8), Inches(4),
        [
            "30-40% del consumo eléctrico mundial proviene de edificios.",
            "Los gestores ven la factura, pero no la fase ni el equipo culpable.",
            "Las anomalías (consumos fuera de patrón) pasan desapercibidas durante días.",
            "Apagar manualmente cargas no críticas requiere coordinación lenta.",
        ],
        size=18,
    )


def build_solution(prs: Presentation, page: int, dashboard: Path) -> None:
    slide = build_title_slide(
        prs, page,
        "La solución",
        "Un panel B2B que ve, alerta y actúa",
        "Streamlit + telemetría en tiempo real + Modo Eco granular por fase",
    )
    slide.shapes.add_picture(str(dashboard), Inches(0.7), Inches(2.4), width=Inches(8.2))
    add_bullets(
        slide, Inches(9.1), Inches(2.5), Inches(3.7), Inches(4.5),
        [
            "Mapa con todos los edificios.",
            "KPIs vivos: kWh, €/MWh, anomalías.",
            "Curvas por fase con marcadores visuales.",
            "Botón de pánico con efecto inmediato.",
        ],
        size=16,
    )


def build_advantages(prs: Presentation, page: int) -> None:
    slide = build_title_slide(
        prs, page,
        "Ventajas clave",
        "Por qué Energy Hunter mueve la aguja",
        "Diseñado alrededor del ahorro, la accesibilidad y la velocidad de reacción",
    )

    cards = [
        ("Ahorro inmediato", "Hasta −40% de consumo en las fases que el operador elige, sin obras."),
        ("Detección de anomalías", "Regla de domingo (×1.40 sobre baseline) y marcadores en mapa y gráfica."),
        ("Control granular", "Modo Eco fase a fase, no \"todo o nada\". Cero impacto en cargas críticas."),
        ("Telemetría 30 s", "Lectura cada 30 s con back-off automático y estado API siempre visible."),
        ("Accesible WCAG 2.1 AA", "Lectores de pantalla, contraste alto, navegación por teclado."),
        ("Stack abierto", "Python + Streamlit + Flask. Reproducible en local en 2 comandos."),
    ]
    col_w = Inches(4.0)
    row_h = Inches(1.55)
    x0 = Inches(0.7)
    y0 = Inches(2.5)
    for i, (title, desc) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = x0 + col * (col_w + Inches(0.15))
        y = y0 + row * (row_h + Inches(0.2))
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, row_h)
        card.line.color.rgb = LEAF_GREEN
        card.line.width = Pt(1.5)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        # Accent dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.22), Inches(0.22), Inches(0.22))
        dot.line.fill.background()
        dot.fill.solid()
        dot.fill.fore_color.rgb = LIME
        add_textbox(slide, x + Inches(0.55), y + Inches(0.15), col_w - Inches(0.6), Inches(0.4),
                    title, size=16, bold=True, color=DEEP_GREEN)
        add_textbox(slide, x + Inches(0.25), y + Inches(0.6), col_w - Inches(0.4), Inches(0.95),
                    desc, size=12, color=CHARCOAL)


def build_panic(prs: Presentation, page: int, image: Path) -> None:
    slide = build_title_slide(
        prs, page,
        "Demo destacada",
        "Botón de Pánico granular: del clic al ahorro en segundos",
        "Marca las fases, confirma, y el webhook propaga Modo Eco al hardware simulado",
    )
    slide.shapes.add_picture(str(image), Inches(0.7), Inches(2.4), width=Inches(7))
    add_bullets(
        slide, Inches(8.1), Inches(2.5), Inches(4.7), Inches(4.5),
        [
            "1. El operador ve un pico anómalo.",
            "2. Selecciona qué fases pasan a Eco.",
            "3. Confirma con un único botón.",
            "4. Webhook → API → −40% en esas fases.",
            "5. Estado se anuncia por aria-live (a11y).",
            "6. Cargas críticas siguen intactas.",
        ],
        size=15,
    )


def build_architecture(prs: Presentation, page: int, image: Path) -> None:
    slide = build_title_slide(
        prs, page,
        "Cómo funciona",
        "Arquitectura simple, demoable en local",
        "Sin dependencias cloud para el hackathon: 2 comandos y listo",
    )
    slide.shapes.add_picture(str(image), Inches(0.7), Inches(2.3), width=Inches(12))
    add_textbox(
        slide, Inches(0.7), Inches(6.2), Inches(12), Inches(0.6),
        "Streamlit consume el Mock API (Smart Breaker + SIOS), Kiro orquesta el "
        "fan-out a Modo Eco. TDD + property-based tests cubren las 10 propiedades críticas.",
        size=13, color=SLATE,
    )


def build_savings(prs: Presentation, page: int, image: Path) -> None:
    slide = build_title_slide(
        prs, page,
        "Impacto",
        "Ahorro medible por categoría de carga",
        "Ejemplo de un edificio piloto: 4 fases activadas en Modo Eco durante 1 hora",
    )
    slide.shapes.add_picture(str(image), Inches(0.7), Inches(2.3), width=Inches(8.5))
    add_bullets(
        slide, Inches(9.4), Inches(2.5), Inches(3.4), Inches(4.5),
        [
            "−40% por fase activada.",
            "Reducción proporcional de €/h.",
            "Menos huella de CO₂ inmediata.",
            "Sin tocar servidores críticos.",
            "Auditable: cada acción queda registrada.",
        ],
        size=15,
    )


def build_aws_deployment(prs: Presentation, page: int, image: Path) -> None:
    slide = build_title_slide(
        prs, page,
        "Despliegue en producción",
        "Cómo se escala Energy Hunter en AWS con IoT Core",
        "Mismo código del MVP, ahora con grado industrial: seguro, elástico y multi-región",
    )
    slide.shapes.add_picture(str(image), Inches(0.4), Inches(2.25), width=Inches(8.4))

    add_bullets(
        slide, Inches(9.0), Inches(2.4), Inches(4.0), Inches(4.6),
        [
            "AWS IoT Core: MQTT mutuo TLS, Device Shadow por fase.",
            "Rules Engine → Lambda → Timestream para series temporales.",
            "S3 + Athena para histórico y auditoría barata.",
            "Lambda detecta anomalías y publica al webhook de Modo Eco.",
            "Streamlit corre en ECS Fargate detrás de CloudFront + Cognito.",
            "Greengrass opcional para autonomía sin red en planta.",
        ],
        size=14,
    )


def build_aws_costs(prs: Presentation, page: int, image: Path) -> None:
    slide = build_title_slide(
        prs, page,
        "Coste y retorno",
        "Lo que paga el cliente en AWS frente a lo que ahorra",
        "Tres escenarios reales · ROI positivo desde el primer mes",
    )
    slide.shapes.add_picture(str(image), Inches(0.4), Inches(2.2), width=Inches(8.4))

    add_bullets(
        slide, Inches(9.0), Inches(2.3), Inches(4.0), Inches(4.7),
        [
            "Pyme: ~95 €/mes en AWS · ahorro ≈ 720 €/mes.",
            "Mediana: ~380 €/mes en AWS · ahorro ≈ 4.200 €/mes.",
            "Corporativo: ~1.850 €/mes · ahorro ≈ 26.000 €/mes.",
            "Pago por uso: IoT Core, Timestream, Lambda, Fargate.",
            "Sin licencias por dispositivo ni costes ocultos.",
            "ROI > 7× en todos los tramos analizados.",
        ],
        size=14,
    )

    # Disclaimer footer
    add_textbox(
        slide, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.35),
        "Estimaciones orientativas en eu-west-1, basadas en mensajes IoT cada 30 s "
        "y ahorro medio del 15% sobre la factura mensual del cliente.",
        size=10, color=SLATE,
    )


def build_aws_breakdown(prs: Presentation, page: int, image: Path) -> None:
    slide = build_title_slide(
        prs, page,
        "Coste — desglose",
        "A dónde va cada euro que pagas a AWS",
        "Sin sorpresas: 5 servicios, todos con tarifa pública y pago por uso",
    )
    slide.shapes.add_picture(str(image), Inches(0.4), Inches(2.2), width=Inches(8.4))
    add_bullets(
        slide, Inches(9.0), Inches(2.3), Inches(4.0), Inches(4.7),
        [
            "IoT Core: 1 $/M mensajes MQTT.",
            "Timestream: writes + memory store + queries.",
            "Lambda: ms de cómputo por anomalía.",
            "Fargate: vCPU/GB-h del Streamlit.",
            "CloudFront, Cognito y S3: tráfico y auth.",
            "Reservas y Savings Plans bajan 20-40%.",
        ],
        size=14,
    )
    add_textbox(
        slide, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.35),
        "Cifras orientativas. La calculadora oficial de AWS las afina por carga real.",
        size=10, color=SLATE,
    )


def build_roadmap(prs: Presentation, page: int, image: Path) -> None:
    slide = build_title_slide(
        prs, page,
        "Hacia dónde va",
        "Roadmap del MVP al producto",
        "Hoy: hackathon. Mañana: piloto real con breakers físicos.",
    )
    slide.shapes.add_picture(str(image), Inches(0.5), Inches(2.5), width=Inches(12.3))


def build_closing(prs: Presentation, page: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DEEP_GREEN)
    add_textbox(
        slide, Inches(0.7), Inches(1.6), Inches(12), Inches(1),
        "Energía bajo control. Naturaleza protegida.",
        size=44, bold=True, color=WHITE,
    )
    add_textbox(
        slide, Inches(0.7), Inches(2.7), Inches(12), Inches(0.8),
        "Energy Hunter convierte cada edificio en un actor activo del ahorro.",
        size=22, color=SAND,
    )

    # Three stat tiles
    stats = [("−40%", "consumo por fase\nen Modo Eco"),
             ("30 s", "frecuencia de\ntelemetría"),
             ("WCAG 2.1 AA", "accesibilidad\npor diseño")]
    for i, (big, small) in enumerate(stats):
        x = Inches(0.7) + i * Inches(4.1)
        tile = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4), Inches(3.9), Inches(2))
        tile.line.fill.background()
        tile.fill.solid()
        tile.fill.fore_color.rgb = LEAF_GREEN
        add_textbox(slide, x, Inches(4.1), Inches(3.9), Inches(0.8),
                    big, size=40, bold=True, color=LIME, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(5.0), Inches(3.9), Inches(1),
                    small, size=14, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(
        slide, Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
        "¿Preguntas? — github.com/energy-hunter   ·   demo en vivo disponible",
        size=14, color=SAND, align=PP_ALIGN.CENTER,
    )


# --------------------------------------------------------------------------- #
# Main                                                                        #
# --------------------------------------------------------------------------- #

def main() -> Path:
    # 1. Generate images
    hero = ASSETS / "hero.png"
    problem = ASSETS / "problem_meter.png"
    dashboard = ASSETS / "dashboard_mock.png"
    panic = ASSETS / "panic_panel.png"
    architecture = ASSETS / "architecture.png"
    aws_arch = ASSETS / "aws_architecture.png"
    cost_compare = ASSETS / "aws_cost_comparison.png"
    cost_breakdown = ASSETS / "aws_cost_breakdown.png"
    savings = ASSETS / "savings_chart.png"
    roadmap = ASSETS / "roadmap.png"

    make_hero_image(hero)
    make_problem_image(problem)
    make_dashboard_mock(dashboard)
    make_panic_image(panic)
    make_architecture_image(architecture)
    make_aws_architecture_image(aws_arch)
    make_cost_comparison_image(cost_compare)
    make_aws_cost_breakdown_image(cost_breakdown)
    make_savings_chart(savings)
    make_roadmap_image(roadmap)

    # 2. Build deck (16:9)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Speaker notes (~25-30 s each, kept tight for a 5-minute pitch).
    notes = {
        1: (
            "Hola, somos el equipo de Energy Hunter. En los próximos cinco minutos vais "
            "a ver una herramienta B2B que ayuda a edificios y empresas a consumir menos "
            "energía y emitir menos CO2, con un solo clic. Menos consumo, menos "
            "emisiones, más control."
        ),
        2: (
            "Los edificios son responsables de hasta el 40 por ciento del consumo "
            "eléctrico mundial. El gestor recibe la factura a fin de mes pero no sabe "
            "qué fase ni qué equipo se le ha disparado. Las anomalías pasan días sin "
            "detectarse, y apagar cargas no críticas implica llamadas y coordinación. "
            "Hay un agujero entre lo que se mide y lo que se puede accionar."
        ),
        3: (
            "Energy Hunter cierra ese hueco. Es un panel hecho en Streamlit que recibe "
            "telemetría cada 30 segundos por fase, muestra los edificios en mapa, "
            "calcula el coste con el precio del mercado SIOS y marca cualquier consumo "
            "anómalo. Ve, alerta, y en la siguiente diapositiva veréis cómo actúa."
        ),
        4: (
            "Estas son las seis ventajas que diferencian al producto: ahorro inmediato "
            "del 40 por ciento por fase, detección automática de anomalías de domingo, "
            "control granular sin tocar cargas críticas, telemetría cada 30 segundos "
            "con back-off, accesibilidad WCAG 2.1 AA y un stack abierto que cualquiera "
            "puede levantar en local en dos comandos."
        ),
        5: (
            "Esta es la demo estrella: el botón de pánico granular. El operador "
            "selecciona qué fases pasan a Modo Eco, confirma, y el webhook propaga la "
            "orden al breaker. Servidores críticos siguen intactos. Iluminación, "
            "climatización y cargadores se reducen un 40 por ciento al instante. La "
            "acción se anuncia con aria-live para lectores de pantalla."
        ),
        6: (
            "La arquitectura del MVP es simple a propósito. Streamlit consume un Mock "
            "API local que simula los breakers y los precios SIOS, y un stub de Kiro "
            "orquesta el fan-out a Modo Eco. Todo se levanta con dos comandos. La "
            "lógica está cubierta con TDD y diez propiedades verificadas con Hypothesis."
        ),
        7: (
            "Para producción, llevamos exactamente la misma lógica a AWS. Los breakers "
            "publican por MQTT a IoT Core con TLS mutuo. Las reglas disparan funciones "
            "Lambda que escriben en Timestream y, si detectan una anomalía, llaman al "
            "endpoint de Modo Eco. El frontend Streamlit corre en Fargate detrás de "
            "CloudFront y Cognito. Greengrass es opcional para autonomía sin red."
        ),
        8: (
            "Y aquí está el argumento de venta: en los tres tramos el ahorro supera al "
            "coste por más de siete a uno. Una pyme paga noventa y cinco euros al mes "
            "en AWS y se ahorra setecientos veinte. Una mediana paga trescientos "
            "ochenta y se ahorra cuatro mil doscientos. Un cliente corporativo paga "
            "mil ochocientos cincuenta y se ahorra veintiséis mil. ROI positivo desde "
            "el primer mes."
        ),
        9: (
            "Si os preguntáis a dónde va cada euro: Timestream y Fargate son los "
            "servicios más pesados, IoT Core suma poco gracias a la tarifa por millón "
            "de mensajes, y Lambda apenas pesa. No hay licencias por dispositivo ni "
            "compromisos anuales. Con Savings Plans bajamos otro veinte a cuarenta por "
            "ciento. Cero sorpresas en factura."
        ),
        10: (
            "El impacto es medible y auditable. Cada fase que entra en Modo Eco baja "
            "su consumo un cuarenta por ciento, lo que se traduce en menos kilovatios, "
            "menos euros y menos CO2 emitido. Y todo queda registrado, así que el "
            "cliente puede demostrar el ahorro a su comité o a su auditor de "
            "sostenibilidad."
        ),
        11: (
            "El roadmap es claro. Hoy entregamos el MVP del hackathon. El siguiente "
            "paso es un piloto con tres edificios reales. Después integramos breakers "
            "físicos vía IoT Core, añadimos multi-tenant y alertas con IA, y lo "
            "llevamos al mercado europeo."
        ),
        12: (
            "Frente a otros proyectos del hackathon, no nos quedamos en la "
            "visualización: actuamos. Granularidad por fase, stack abierto, "
            "accesibilidad desde el día uno y verificación formal con propiedades. "
            "Eso convierte una demo bonita en algo que un comprador B2B puede llevar "
            "a producción."
        ),
        13: (
            "Energía bajo control, naturaleza protegida. Energy Hunter convierte cada "
            "edificio en un actor activo del ahorro. Estamos abiertos a preguntas y "
            "tenemos demo en vivo lista. Gracias."
        ),
    }

    s1 = build_cover(prs, hero)                                           # 1
    s2 = build_problem(prs, 2, problem)                                   # 2
    s3 = build_solution(prs, 3, dashboard)                                # 3
    s4 = build_advantages(prs, 4)                                         # 4
    s5 = build_panic(prs, 5, panic)                                       # 5
    s6 = build_architecture(prs, 6, architecture)                         # 6
    s7 = build_aws_deployment(prs, 7, aws_arch)                           # 7
    s8 = build_aws_costs(prs, 8, cost_compare)                            # 8
    s9 = build_aws_breakdown(prs, 9, cost_breakdown)                      # 9
    s10 = build_savings(prs, 10, savings)                                 # 10
    s11 = build_roadmap(prs, 11, roadmap)                                 # 11
    # Slide 12: differentiation vs alternatives
    s12 = build_title_slide(
        prs, 12,
        "Frente a alternativas",
        "Lo que nos diferencia de otros proyectos del hackathon",
        "Foco en acción, no solo en visualización",
    )
    add_bullets(
        s12, Inches(0.7), Inches(2.6), Inches(12), Inches(4),
        [
            "Otros muestran datos: nosotros además accionamos el Modo Eco en segundos.",
            "Granularidad por fase, no por edificio entero (cargas críticas siempre seguras).",
            "Stack 100% open-source y reproducible en local sin cuentas cloud.",
            "Accesibilidad WCAG AA desde el día 1, no como capa cosmética.",
            "TDD con propiedades formales: 10 invariantes verificadas con Hypothesis.",
        ],
        size=18,
    )
    s13 = build_closing(prs, 13)                                          # 13

    # Attach speaker notes (slides are 1-indexed via prs.slides)
    slide_objs = list(prs.slides)
    for idx, slide_obj in enumerate(slide_objs, start=1):
        if idx in notes:
            add_speaker_notes(slide_obj, notes[idx])

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    out = main()
    print(f"Deck escrito en: {out}")
